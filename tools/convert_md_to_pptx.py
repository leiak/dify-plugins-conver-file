from collections.abc import Generator
from typing import Any
import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import re

from dify_plugin import Tool
from dify_plugin.entities.tool import ToolInvokeMessage

class ConvertMdToPptxTool(Tool):
    def _invoke(self, tool_parameters: dict[str, Any]) -> Generator[ToolInvokeMessage, None, None]:
        content = tool_parameters.get("content", "")
        filename = tool_parameters.get("filename", "output")
        
        if not content:
            raise Exception("Content cannot be empty.")

        # 创建演示文稿
        prs = Presentation()
        
        # 解析 Markdown 内容
        lines = content.split('\n')
        
        # 状态变量
        current_title = "Presentation"
        current_slide = None
        
        # 添加标题和内容幻灯片的辅助函数
        def add_content_slide(title_text, body_text):
            layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            
            # Set content
            if slide.placeholders and len(slide.placeholders) > 1:
                content_shape = slide.placeholders[1]
                content_shape.text = body_text
                
                # Adjust font size if needed (optional, simple heuristic)
                if len(body_text) > 500:
                    for paragraph in content_shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
            
            return slide

        # 预处理行，将其分组为块（标题、段落）
        blocks = []
        current_block = []
        
        for line in lines:
            line = line.strip()
            if not line:
                if current_block:
                    blocks.append(current_block)
                    current_block = []
                continue
            
            if line.startswith('#'):
                if current_block:
                    blocks.append(current_block)
                    current_block = []
                blocks.append([line])
            else:
                current_block.append(line)
        
        if current_block:
            blocks.append(current_block)

        # 根据块生成幻灯片
        for block in blocks:
            first_line = block[0]
            
            # 处理 H1 - 标题幻灯片
            if first_line.startswith('# '):
                title_text = first_line[2:].strip()
                layout = prs.slide_layouts[0] # Title Slide
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = title_text
                current_title = title_text # Update context title
                
            # 处理 H2 - 新章节开始
            elif first_line.startswith('## '):
                title_text = first_line[3:].strip()
                current_title = title_text
                # If there is content in this block (lines after header), add it
                if len(block) > 1:
                    body_text = '\n'.join(block[1:])
                    add_content_slide(current_title, body_text)
                else:
                    # Just a header, create a slide to establish the section
                    add_content_slide(current_title, "")
                    
            # 处理段落/列表
            else:
                # It's a content block. Create a new slide for it using the current title.
                # This ensures each paragraph/block gets its own slide, preventing overcrowding.
                body_text = '\n'.join(block)
                add_content_slide(current_title, body_text)

        # 如果未创建幻灯片，则创建默认幻灯片
        if not prs.slides:
            add_content_slide("Converted Content", content)

        # 保存到缓冲区
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_bytes = pptx_buffer.getvalue()
        
        # 将 PPTX 作为 blob 消息返回
        yield self.create_blob_message(
            blob=pptx_bytes, 
            meta={'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'filename': f"{filename}.pptx"}
        )
